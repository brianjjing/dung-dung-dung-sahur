#!/usr/bin/env python3
"""Generate the Triple D pitch deck (.pptx) — clean, professional, white.

Styled after the original AirBnB pitch deck: white background, generous
whitespace, section title top-left + page number top-right in ONE accent
color, sparse text, flat color boxes, big bold numbers. The brand logo
appears ONLY on the first (title) and the closing slide. Content slides are
logo-free. Appendix slides are unbranded backup for Q&A.

Run:  python brand/make_deck.py
Out:  brand/TripleD.pptx   (16:9)

Single accent = Defeat Amber (the brand signature, playing AirBnB-blue's role).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ------------------------------------------------------------- palette (white)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # background
INK   = RGBColor(0x16, 0x21, 0x2B)   # headings emphasis / near-black body
SLATE = RGBColor(0x6B, 0x76, 0x83)   # secondary / continuation text
FAINT = RGBColor(0xC2, 0xC9, 0xD1)   # hairlines, footer
AMBER = RGBColor(0xF0, 0x8A, 0x12)   # THE accent (titles, numbers, boxes)
BLACK = RGBColor(0x0B, 0x0F, 0x14)   # closing-slide background (Command Black)
# brand state colors — used only for the logo mark
CYAN  = RGBColor(0x27, 0xE0, 0xC8)
AMBR  = RGBColor(0xFF, 0x9E, 0x2C)
RED   = RGBColor(0xF0, 0x50, 0x3A)

F_HEAD = "Helvetica Neue"   # titles + body (clean, neutral, professional)
F_BODY = "Helvetica Neue"
F_MONO = "Menlo"            # only for the literal command / serial strings

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------- helpers
def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf


def run(p, text, size, color, bold=False, italic=False, font=F_HEAD, spc=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    if spc is not None:
        r._r.get_or_add_rPr().set('spc', str(int(spc * 100)))
    return r


def para(tf, first=False, align=PP_ALIGN.LEFT, before=0, after=6, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before); p.space_after = Pt(after)
    if line is not None:
        p.line_spacing = line
    return p


def rect(s, x, y, w, h, fill=None, ln=None, ln_w=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if ln is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = ln; sh.line.width = Pt(ln_w)
    sh.shadow.inherit = False
    return sh


def header(s, title_text, num, accent=AMBER):
    """Section title top-left + page number top-right (the AirBnB signature)."""
    tb, tf = box(s, 0.85, 0.5, 9.5, 1.0)
    run(para(tf, first=True), title_text, 38, accent, bold=True)
    tb, tf = box(s, 10.5, 0.5, 1.98, 1.0)
    run(para(tf, first=True, align=PP_ALIGN.RIGHT), str(num), 38, accent,
        bold=True)


def footer(s):
    tb, tf = box(s, 9.5, 7.06, 2.98, 0.35)
    run(para(tf, first=True, align=PP_ALIGN.RIGHT), "TRIPLE D", 9, FAINT,
        bold=True, spc=2)


def logo(s, x, y, scale=1.0, mono=None):
    """Triple-chevron mark, cyan -> amber -> red. mono=color overrides all."""
    size, gap = 0.62 * scale, 0.16 * scale
    cols = [mono or CYAN, mono or AMBR, mono or RED]
    for i, c in enumerate(cols):
        ch = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                Inches(x + i * (size * 0.55 + gap)),
                                Inches(y), Inches(size), Inches(size))
        ch.fill.solid(); ch.fill.fore_color.rgb = c
        ch.line.fill.background(); ch.shadow.inherit = False


def lead(tf, bold_text, rest_text, size=27, p=None, after=18, line=1.2):
    """AirBnB-style statement: bold near-black lead + slate continuation."""
    p = p or para(tf, after=after, line=line)
    run(p, bold_text, size, INK, bold=True, font=F_BODY)
    run(p, rest_text, size, SLATE, font=F_BODY)
    return p


# ============================================================= 1 · TITLE (logo)
s = slide()
logo(s, 5.55, 2.0, scale=1.15)
tb, tf = box(s, 1.0, 3.15, 11.333, 1.4)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "TRIPLE ", 76, INK, bold=True)
run(p, "D", 76, AMBER, bold=True)
tb, tf = box(s, 1.0, 4.55, 11.333, 0.7)
run(para(tf, first=True, align=PP_ALIGN.CENTER),
    "Detect.  Decide.  Defeat.", 26, SLATE, font=F_BODY)
tb, tf = box(s, 1.0, 5.4, 11.333, 0.6)
run(para(tf, first=True, align=PP_ALIGN.CENTER),
    "Non-kinetic, human-accountable counter-drone defense.", 17, SLATE,
    font=F_BODY)
# subtle bottom-right context note (like AirBnB's grey note box)
tb, tf = box(s, 8.0, 6.55, 4.5, 0.6)
run(para(tf, first=True, align=PP_ALIGN.RIGHT),
    "Against RF-silent, fiber-optic drones.", 13, FAINT, italic=True,
    font=F_BODY)

# ================================================================= 2 · PROBLEM
s = slide()
header(s, "Problem", 2)
tb, tf = box(s, 1.4, 2.15, 10.3, 4.5)
lead(tf, "RF-silent.  ",
     "Fiber-optic suicide drones carry no radio link — jammers and "
     "electronic-warfare systems see nothing.", p=para(tf, first=True,
     after=22, line=1.2))
lead(tf, "Wire-guided.  ",
     "They fly on a hair-thin spool of glass. There's no signal to cut and "
     "no GPS to spoof.", after=22)
lead(tf, "Fielded at scale.  ",
     "Cheap, precise, and increasingly common — with a human on the other "
     "end making a kill decision.", after=0)
footer(s)

# ================================================================ 3 · SOLUTION
s = slide()
header(s, "Solution", 3)
tb, tf = box(s, 1.4, 2.05, 10.5, 1.1)
lead(tf, "A layered, non-kinetic kill-chain  ",
     "that turns sound into a decision, and a decision into a safe response.",
     size=24, p=para(tf, first=True, line=1.2))
boxes = [
    ("DETECT", "acoustic gate"),
    ("DECIDE", "vision + fusion"),
    ("DEFEAT", "layered response"),
]
bw, bx0, gap = 3.55, 1.4, 0.42
for i, (word, cap) in enumerate(boxes):
    x = bx0 + i * (bw + gap)
    rect(s, x, 3.5, bw, 2.3, fill=AMBER)
    tb, tf = box(s, x, 3.5, bw, 2.3, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER, after=4), word, 30, WHITE,
        bold=True)
    run(para(tf, align=PP_ALIGN.CENTER), cap, 15, WHITE, font=F_BODY)
footer(s)

# ============================================================ 4 · HOW IT WORKS
s = slide()
header(s, "How It Works", 4)
# flow row (SEARCH -> REVIEW -> BOOK archetype)
steps = ["HEAR", "CONFIRM", "AUTHORIZE", "DEFEAT"]
subs  = ["acoustic", "vision", "human", "effect"]
fx, fy, step_w = 1.4, 2.25, 2.35
for i, (st, sub) in enumerate(zip(steps, subs)):
    x = fx + i * step_w
    tb, tf = box(s, x, fy, step_w - 0.4, 0.85)
    run(para(tf, first=True, after=2), st, 22, INK, bold=True)
    run(para(tf), sub, 13, SLATE, font=F_BODY, spc=1)
    if i < len(steps) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(x + step_w - 0.62), Inches(fy + 0.12),
                                Inches(0.5), Inches(0.28))
        ar.fill.solid(); ar.fill.fore_color.rgb = FAINT
        ar.line.fill.background(); ar.shadow.inherit = False
# two brains (light-outline boxes)
rect(s, 1.4, 3.55, 4.55, 2.25, fill=None, ln=FAINT, ln_w=1.25)
tb, tf = box(s, 1.7, 3.8, 4.0, 1.8)
run(para(tf, first=True, after=8), "UNO R3 — on the car", 16, AMBER, bold=True)
for t in ["Thin C firmware; flashed once", "Mic features, ultrasonic distance",
          "Drives motors, fires effects"]:
    run(para(tf, after=4, line=1.1), "·  " + t, 14, SLATE, font=F_BODY)
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.15), Inches(4.4),
                        Inches(0.85), Inches(0.32))
ar.fill.solid(); ar.fill.fore_color.rgb = FAINT
ar.line.fill.background(); ar.shadow.inherit = False
rect(s, 7.25, 3.55, 4.55, 2.25, fill=None, ln=AMBER, ln_w=1.5)
tb, tf = box(s, 7.55, 3.8, 4.0, 1.8)
run(para(tf, first=True, after=8), "UNO Q — the brain", 16, AMBER, bold=True)
for t in ["Python state machine (offline)", "DETECT, DECIDE, fusion",
          "Operator gate, DEFEAT commands"]:
    run(para(tf, after=4, line=1.1), "·  " + t, 14, SLATE, font=F_BODY)
tb, tf = box(s, 1.4, 6.15, 10.5, 0.7)
p = para(tf, first=True, line=1.15)
run(p, "The car never decides anything.  ", 17, INK, bold=True, font=F_BODY)
run(p, "All intelligence runs offline on the Uno Q.", 17, SLATE, font=F_BODY)
footer(s)

# ============================================================= 5 · RESPONSE HEADS
s = slide()
header(s, "Response Heads", 5)
heads = [
    ("DISTRACT", "Decoy LEDs", "Harmless. Confuses the optical seeker and draws aim away. The one action safe to automate."),
    ("DISABLE", "IR dazzle", "Blinds the drone's seeker optics. Always human-authorized."),
    ("DEFEND", "Precision lock & track", "Tracking only — maintains lock on target. No kinetic action, by design."),
]
bw, bx0, gap = 3.55, 1.4, 0.42
for i, (word, sub, body) in enumerate(heads):
    x = bx0 + i * (bw + gap)
    rect(s, x, 2.35, bw, 0.85, fill=AMBER)
    tb, tf = box(s, x, 2.35, bw, 0.85, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), word, 22, WHITE, bold=True)
    tb, tf = box(s, x + 0.05, 3.45, bw - 0.1, 0.5)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), sub, 16, INK, bold=True,
        font=F_BODY)
    tb, tf = box(s, x + 0.1, 4.05, bw - 0.2, 2.0)
    run(para(tf, first=True, align=PP_ALIGN.CENTER, line=1.25), body, 14, SLATE,
        font=F_BODY)
tb, tf = box(s, 1.4, 6.35, 10.5, 0.6)
p = para(tf, first=True)
run(p, "Escalating and layered — ", 15, INK, bold=True, font=F_BODY)
run(p, "the harmless head can run on its own; the rest wait for a human.", 15,
    SLATE, font=F_BODY)
footer(s)

# ====================================================== 6 · RESPONSIBLE AUTONOMY
s = slide()
header(s, "Responsible Autonomy", 6)
tb, tf = box(s, 1.4, 2.05, 10.5, 1.3)
p = para(tf, first=True, line=1.2)
run(p, "The safe action automates.\n", 28, INK, bold=True, font=F_BODY)
run(p, "The harmful one always asks a human.", 28, AMBER, bold=True,
    font=F_BODY)
rows = [
    ("L0", "Teleop", "human does everything"),
    ("L2", "Detect + recommend", "human gates all actions  ·  demo default"),
    ("L3", "Human-on-the-loop", "system acts; human can veto in a window"),
    ("L4", "Bounded autonomy", "auto-DISTRACT only; gate DISABLE / DEFEND"),
]
y = 3.85
for lab, name, desc in rows:
    highlight = lab == "L2"
    rect(s, 1.4, y, 0.85, 0.6, fill=(AMBER if highlight else None),
         ln=(None if highlight else FAINT), ln_w=1.25)
    tb, tf = box(s, 1.4, y, 0.85, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), lab, 18,
        (WHITE if highlight else AMBER), bold=True)
    tb, tf = box(s, 2.55, y, 3.6, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True), name, 17, INK, bold=True, font=F_BODY)
    tb, tf = box(s, 6.3, y, 6.2, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True), desc, 15, SLATE, font=F_BODY)
    y += 0.72
footer(s)

# ================================================================== 7 · SCOPE
s = slide()
header(s, "Scope", 7)
tb, tf = box(s, 1.4, 2.1, 10.5, 0.7)
lead(tf, "Scoped on purpose.  ",
     "Credibility comes from naming the edges — the demo does exactly this, "
     "no more.", size=21, p=para(tf, first=True, line=1.2))
items = [
    ("DEFEND is tracking-only.", "Cutting a hair-thin moving fiber is out of scope. The demo is precision lock-and-track; no “cut” command exists."),
    ("DISABLE defeats the seeker.", "IR dazzle blinds the optics — not “disarm the warhead.”"),
    ("One mic = detection, not bearing.", "Direction-finding needs multiple nodes. We detect presence, not azimuth."),
    ("L0 teleop is always the fallback.", "Every demo is rehearsed with a human able to take full manual control."),
]
y = 3.05
for head_t, body in items:
    rect(s, 1.4, y + 0.07, 0.16, 0.16, fill=AMBER)
    tb, tf = box(s, 1.85, y, 10.4, 0.9)
    p = para(tf, first=True, line=1.15)
    run(p, head_t + "   ", 18, INK, bold=True, font=F_BODY)
    run(p, body, 16, SLATE, font=F_BODY)
    y += 0.98
footer(s)

# =========================================================== 8 · CLOSING (logo)
s = slide(bg=BLACK)
logo(s, 5.55, 1.9, scale=1.0)            # gradient mark pops on dark
tb, tf = box(s, 1.0, 3.05, 11.333, 1.0)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "See it ", 40, WHITE, bold=True)
run(p, "Detect › Decide › Defeat", 40, AMBR, bold=True)
tb, tf = box(s, 1.0, 4.15, 11.333, 0.6)
run(para(tf, first=True, align=PP_ALIGN.CENTER),
    "Prop-whine fires DETECT · vision confirms HOSTILE · operator authorizes · "
    "DEFEAT actuates.", 16, FAINT, font=F_BODY)
rect(s, 4.4, 5.0, 4.5, 0.62, fill=RGBColor(0x16, 0x1C, 0x24), ln=AMBR,
     ln_w=1.0)
tb, tf = box(s, 4.65, 5.0, 4.0, 0.62, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True)
run(p, "$ ", 16, CYAN, bold=True, font=F_MONO)
run(p, "cd uno_q && python main.py", 16, WHITE, font=F_MONO)
tb, tf = box(s, 1.0, 6.25, 11.333, 0.5)
run(para(tf, first=True, align=PP_ALIGN.CENTER),
    "TRIPLE D  —  Non-kinetic defense, human-accountable.", 14, AMBR,
    italic=True, font=F_BODY, spc=1)

# ============================================== APPENDIX (unbranded, Q&A backup)
def appendix(title_text, code):
    s = slide()
    header(s, title_text, code, accent=SLATE)
    return s

# A1 — architecture / serial protocol
s = appendix("Appendix — Architecture", "A1")
tb, tf = box(s, 1.4, 2.05, 10.5, 0.6)
run(para(tf, first=True), "Human-readable serial protocol — open a monitor and "
    "read the whole conversation live.", 16, SLATE, font=F_BODY)
rect(s, 1.4, 2.85, 10.4, 1.35, fill=RGBColor(0xF4, 0xF6, 0xF8), ln=FAINT)
tb, tf = box(s, 1.7, 3.05, 9.9, 1.0)
p = para(tf, first=True)
run(p, "UP    car → Uno Q    ", 14, AMBER, bold=True, font=F_MONO)
run(p, "TEL,AMP:512,PITCH:2200,DIST:84,LINE:0", 14, INK, font=F_MONO)
p = para(tf, line=1.4)
run(p, "DOWN  Uno Q → car    ", 14, AMBER, bold=True, font=F_MONO)
run(p, "CMD,DISTRACT_ON   (DAZZLE, DRIVE_*, ALL_OFF, IDLE)", 14, INK,
    font=F_MONO)
tb, tf = box(s, 1.4, 4.55, 10, 0.45)
run(para(tf, first=True), "STATE MACHINE", 12, SLATE, bold=True, spc=2)
states = ["IDLE", "DECIDING", "AUTHORIZING", "DEFEATING", "COOLDOWN"]
x = 1.4
for i, st in enumerate(states):
    w = 2.0
    rect(s, x, 5.05, w, 0.6, fill=None, ln=(AMBER if st == "AUTHORIZING" else FAINT),
         ln_w=(1.5 if st == "AUTHORIZING" else 1.0))
    tb, tf = box(s, x, 5.05, w, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), st, 13,
        (AMBER if st == "AUTHORIZING" else INK), bold=True, font=F_MONO)
    x += w + 0.1
footer(s)

# A2 — autonomy dial (full L0-L5)
s = appendix("Appendix — Autonomy Dial", "A2")
dial = [
    ("L0", "Teleop", "human does everything"),
    ("L1", "Single-action assist", "human triggers each effect"),
    ("L2", "Detect + recommend", "human gates all actions  ·  demo default"),
    ("L3", "Human-on-the-loop", "system acts; human can veto in a window"),
    ("L4", "Bounded autonomy", "auto-DISTRACT only; gate DISABLE / DEFEND"),
    ("L5", "Multi-agent", "swarm hands off targets over the mesh"),
]
y = 2.1
for lab, name, desc in dial:
    hl = lab == "L2"
    rect(s, 1.4, y, 0.8, 0.6, fill=(AMBER if hl else None),
         ln=(None if hl else FAINT), ln_w=1.0)
    tb, tf = box(s, 1.4, y, 0.8, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), lab, 16,
        (WHITE if hl else AMBER), bold=True)
    tb, tf = box(s, 2.5, y, 3.7, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True), name, 16, INK, bold=True, font=F_BODY)
    tb, tf = box(s, 6.3, y, 6.2, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True), desc, 14, SLATE, font=F_BODY)
    y += 0.78
footer(s)

# A3 — tech stack
s = appendix("Appendix — Tech Stack", "A3")
tb, tf = box(s, 1.4, 2.05, 10.5, 0.5)
run(para(tf, first=True), "100% offline, on-device.", 18, INK, bold=True,
    font=F_BODY)
cols = [
    ("DETECT", ["Mel-spectrogram CNN", "PyTorch (.pth)", "Sliding-window inference", "Car-mic acoustic features"]),
    ("DECIDE", ["YOLOv8 ONNX detector", "Roboflow drone dataset", "Closing-behavior fusion", "Vote-window debounce"]),
    ("PLATFORM", ["Arduino Uno Q (Linux)", "Uno R3 C firmware", "USB serial @ 115200", "JSONL audit log"]),
]
bw, bx0, gap = 3.55, 1.4, 0.42
for i, (name, items) in enumerate(cols):
    x = bx0 + i * (bw + gap)
    rect(s, x, 2.85, bw, 0.7, fill=AMBER)
    tb, tf = box(s, x, 2.85, bw, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), name, 18, WHITE, bold=True)
    tb, tf = box(s, x + 0.2, 3.8, bw - 0.3, 2.5)
    for it in items:
        run(para(tf, after=8, line=1.15), "·  " + it, 14, SLATE, font=F_BODY)
footer(s)

prs.save("brand/TripleD.pptx")
print(f"Saved brand/TripleD.pptx — {len(prs.slides._sldIdLst)} slides")
